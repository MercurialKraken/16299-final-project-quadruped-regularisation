"""Add discordant push video + rough terrain slides to the deck.
Inserts before the conclusion slide.
"""
import io, sys, os
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

SRC = r"A:\AllIsaac\IsaacLab\flow_matching_presentation.pptx"
DST = r"A:\AllIsaac\IsaacLab\flow_matching_presentation_v2.pptx"
ROUGH_PLOTS = r"A:\AllIsaac\IsaacLab\rough_plots"
PUSH_VID_DIR = r"A:\AllIsaac\IsaacLab\push_videos"

BG      = RGBColor(0x0F, 0x17, 0x2A)
CARD    = RGBColor(0x1E, 0x29, 0x3B)
TEXT    = RGBColor(0xF8, 0xFA, 0xFC)
MUTED   = RGBColor(0x94, 0xA3, 0xB8)
ACCENT  = RGBColor(0x0E, 0xA5, 0xE9)
ACCENT2 = RGBColor(0x38, 0xBD, 0xF8)
GREEN   = RGBColor(0x22, 0xC5, 0x5E)
RED     = RGBColor(0xEF, 0x44, 0x44)
SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)

def set_slide_bg(slide, rgb):
    bg = slide.background; fill = bg.fill; fill.solid(); fill.fore_color.rgb = rgb

def add_rect(slide, x, y, w, h, fill_rgb):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill_rgb
    shp.line.fill.background(); shp.shadow.inherit = False
    return shp

def add_text(slide, x, y, w, h, text, size=14, color=TEXT, bold=False, font="Calibri", align="left"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
    run = p.add_run(); run.text = text
    run.font.name = font; run.font.size = Pt(size)
    run.font.bold = bold; run.font.color.rgb = color
    return tb

def add_stat_card(slide, x, y, w, h, value, label, value_color=GREEN, value_size=28):
    add_rect(slide, x, y, w, h, CARD)
    add_text(slide, x, y + Emu(int(h * 0.10)), w, Emu(int(h * 0.55)),
             value, size=value_size, color=value_color, bold=True, font="Georgia", align="center")
    add_text(slide, x, y + Emu(int(h * 0.68)), w, Emu(int(h * 0.28)),
             label, size=10, color=MUTED, font="Calibri", align="center")

def add_title_block(slide, title, subtitle):
    add_text(slide, Inches(0.5), Inches(0.25), Inches(9.0), Inches(0.55),
             title, size=22, color=TEXT, bold=True, font="Georgia")
    add_text(slide, Inches(0.5), Inches(0.85), Inches(9.0), Inches(0.30),
             subtitle, size=12, color=MUTED, font="Calibri")
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Emu(55000))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()

def add_footer(slide, text):
    add_text(slide, Inches(0.5), Inches(5.35), Inches(9.0), Inches(0.25),
             text, size=9, color=MUTED, font="Calibri", align="right")

def move_slide(prs, old_index, new_index):
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[old_index])
    xml_slides.insert(new_index, slides[old_index])

def find_mp4(folder):
    if not os.path.isdir(folder): return None
    for f in os.listdir(folder):
        if f.lower().endswith(".mp4"): return os.path.join(folder, f)
    return None

# ============================================================================
# MAIN
# ============================================================================
prs = Presentation(SRC)
n_before = len(prs.slides)
print(f"Original slides: {n_before}")
conclusion_idx = n_before - 1  # conclusion is last
blank = prs.slide_layouts[0]
slides_added = 0

# ============================================================================
# Slide A: Discordant Push Recovery Video (200 N -- raw falls, smoothed survives)
# ============================================================================
raw_mp4 = find_mp4(os.path.join(PUSH_VID_DIR, "disc_raw_200"))
smo_mp4 = find_mp4(os.path.join(PUSH_VID_DIR, "disc_smo_200"))

if raw_mp4 and smo_mp4:
    sa = prs.slides.add_slide(blank)
    set_slide_bg(sa, BG)
    add_title_block(sa,
        "Push Recovery: Discordant Outcome  (F = 200 N)",
        "Same seed, same 200 N lateral impulse  \u2014  raw PPO falls, flow-smoothed survives",
    )

    vid_y = Inches(1.55)
    vid_w = Inches(4.50); vid_h = Inches(2.85); label_h = Inches(0.40)
    gap_x = Inches(0.30)
    total_w = 2 * vid_w + gap_x
    start_x = (SLIDE_W - total_w) / 2

    lx = start_x
    add_text(sa, lx, vid_y, vid_w, label_h, "Raw PPO  \u2014  FELL at step 130",
             size=16, color=RED, bold=True, font="Georgia", align="center")
    sa.shapes.add_movie(raw_mp4, lx, vid_y + label_h, vid_w, vid_h, mime_type="video/mp4")

    rx = start_x + vid_w + gap_x
    add_text(sa, rx, vid_y, vid_w, label_h, "Flow-Smoothed  \u2014  SURVIVED",
             size=16, color=GREEN, bold=True, font="Georgia", align="center")
    sa.shapes.add_movie(smo_mp4, rx, vid_y + label_h, vid_w, vid_h, mime_type="video/mp4")

    add_text(sa, Inches(0.5), Inches(5.00), Inches(9.0), Inches(0.35),
             "Click each video to play. Lateral 200 N impulse for 200 ms at t = 1.6 s. "
             "The smoothed policy's reduced jitter allows more stable recovery from perturbation.",
             size=10, color=MUTED, font="Calibri", align="center")
    add_footer(sa, "push_video_eval.py | seed=42 | F=200 N | push_dur=10 steps (200 ms)")
    slides_added += 1
    print("Added discordant push video slide")
else:
    print(f"WARNING: Discordant videos not found: raw={raw_mp4}, smo={smo_mp4}")

# ============================================================================
# Slide B: Rough Terrain - Action Smoothness Comparison
# ============================================================================
sb = prs.slides.add_slide(blank)
set_slide_bg(sb, BG)
add_title_block(sb,
    "Rough Terrain: Action Smoothness",
    "Flow matching on Isaac-Velocity-Rough-Unitree-Go1-v0 (235-dim obs, height scan)",
)

# Stat row
stat_y = Inches(1.40); stat_h = Inches(0.95)
card_w = Inches(2.85); gap_x = Inches(0.22)
total_w = 3 * card_w + 2 * gap_x
start_x = (SLIDE_W - total_w) / 2

add_stat_card(sb, start_x + 0*(card_w+gap_x), stat_y, card_w, stat_h,
              "15.0%", "Action jitter reduction (t=0.5)", GREEN)
add_stat_card(sb, start_x + 1*(card_w+gap_x), stat_y, card_w, stat_h,
              "24.9%", "Jitter reduction (t=1.0)", GREEN)
add_stat_card(sb, start_x + 2*(card_w+gap_x), stat_y, card_w, stat_h,
              "235-dim", "Observation space (+ height scan)", ACCENT2)

# Bottom: bar chart (left) + action traces (right)
plot_y = Inches(2.50); plot_h = Inches(2.65)
left_x = Inches(0.35); right_x = Inches(5.25); plot_w = Inches(4.50)

add_rect(sb, left_x - Inches(0.05), plot_y - Inches(0.05), plot_w + Inches(0.10), plot_h + Inches(0.10), CARD)
sb.shapes.add_picture(rf"{ROUGH_PLOTS}\rough_smoothness_bars.png", left_x, plot_y, width=plot_w)

add_rect(sb, right_x - Inches(0.05), plot_y - Inches(0.05), plot_w + Inches(0.10), plot_h + Inches(0.10), CARD)
sb.shapes.add_picture(rf"{ROUGH_PLOTS}\rough_action_traces.png", right_x, plot_y, width=plot_w)

add_footer(sb, "5 episodes, 3805 timesteps | flow_model_rough.pt | state_dim=235")
slides_added += 1
print("Added rough terrain smoothness slide")

# ============================================================================
# Slide C: Rough Terrain - Velocity Tracking & Summary
# ============================================================================
sc = prs.slides.add_slide(blank)
set_slide_bg(sc, BG)
add_title_block(sc,
    "Rough Terrain: Velocity Tracking & Summary",
    "Flow smoothing reduces jitter without degrading command tracking on uneven ground",
)

# Velocity tracking plot (left) + summary table (right)
plot_y = Inches(1.50); plot_h = Inches(3.20)
left_x = Inches(0.35); right_x = Inches(5.25); plot_w = Inches(4.50)

add_rect(sc, left_x - Inches(0.05), plot_y - Inches(0.05), plot_w + Inches(0.10), plot_h + Inches(0.10), CARD)
sc.shapes.add_picture(rf"{ROUGH_PLOTS}\rough_velocity_tracking.png", left_x, plot_y, width=plot_w)

add_rect(sc, right_x - Inches(0.05), plot_y - Inches(0.05), plot_w + Inches(0.10), plot_h + Inches(0.10), CARD)
sc.shapes.add_picture(rf"{ROUGH_PLOTS}\rough_summary_table.png", right_x, plot_y, width=plot_w)

add_text(sc, Inches(0.5), Inches(4.90), Inches(9.0), Inches(0.45),
         "Rough terrain preserves less jitter reduction (15-25%) vs flat (65%) because "
         "actions carry more legitimate terrain-response signal. The flow model correctly "
         "distinguishes meaningful high-freq content from noise.",
         size=11, color=TEXT, font="Calibri", align="center")
add_footer(sc, "PPO checkpoint: model_100.pt (100 iters) | rough terrain curriculum")
slides_added += 1
print("Added rough terrain summary slide")

# ============================================================================
# Move conclusion to end
# ============================================================================
final_count = len(prs.slides)
# Conclusion was at conclusion_idx; slides were appended at the end.
# We need to move conclusion (still at its original index) to the very end.
move_slide(prs, conclusion_idx, final_count - 1)

print(f"Final slide count: {final_count} ({slides_added} added)")
prs.save(DST)
print(f"Saved to {DST}")
