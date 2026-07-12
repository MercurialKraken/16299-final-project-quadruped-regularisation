"""Insert effort/energy + push-recovery slides into the deck.

Layout matches add_ab_slides.py (same dark theme).
Inserts 3 new slides:
    1. Effort & energy — bar chart + paired scatter
    2. Effort over time — power vs time
    3. Push recovery — fall rate + McNemar discordant
All three go BEFORE the original conclusion slide.
"""
import copy, io, sys, os
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

SRC = r"A:\IsaacLab\flow_matching_presentation.pptx"
DST = r"A:\IsaacLab\flow_matching_presentation.pptx"
PLOTS = r"A:\IsaacLab\eff_push_plots"
PUSH_VID_DIR = r"A:\IsaacLab\push_videos"

BG      = RGBColor(0x0F, 0x17, 0x2A)
CARD    = RGBColor(0x1E, 0x29, 0x3B)
TEXT    = RGBColor(0xF8, 0xFA, 0xFC)
MUTED   = RGBColor(0x94, 0xA3, 0xB8)
ACCENT  = RGBColor(0x0E, 0xA5, 0xE9)
ACCENT2 = RGBColor(0x38, 0xBD, 0xF8)
GREEN   = RGBColor(0x22, 0xC5, 0x5E)
RED     = RGBColor(0xEF, 0x44, 0x44)

SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)


def set_slide_bg(slide, rgb):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def add_rect(slide, x, y, w, h, fill_rgb, line_rgb=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_rgb
    if line_rgb is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_rgb
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, size=14, color=TEXT, bold=False, font="Calibri", align="left"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_stat_card(slide, x, y, w, h, value, label, value_color=GREEN, value_size=28):
    add_rect(slide, x, y, w, h, CARD)
    add_text(slide, x, y + Emu(int(h * 0.10)), w, Emu(int(h * 0.55)),
             value, size=value_size, color=value_color, bold=True, font="Georgia", align="center")
    add_text(slide, x, y + Emu(int(h * 0.68)), w, Emu(int(h * 0.28)),
             label, size=10, color=MUTED, font="Calibri", align="center")


def add_title_block(slide, title, subtitle):
    add_text(slide, Inches(0.5), Inches(0.25), Inches(9.0), Inches(0.55),
             title, size=22, color=TEXT, bold=True, font="Georgia")
    add_text(slide, Inches(0.5), Inches(0.85), Inches(9.0), Inches(0.30),
             subtitle, size=12, color=MUTED, font="Calibri")
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Emu(55000))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()


def add_footer(slide, text):
    add_text(slide, Inches(0.5), Inches(5.35), Inches(9.0), Inches(0.25),
             text, size=9, color=MUTED, font="Calibri", align="right")


def move_slide(prs, old_index, new_index):
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[old_index])
    xml_slides.insert(new_index, slides[old_index])


prs = Presentation(SRC)
n_before = len(prs.slides)
print(f"Original slides: {n_before}")

# Conclusion is currently last slide
conclusion_old_index = n_before - 1
blank = prs.slide_layouts[0]


# ============================================================================
# Slide 1: Control effort & energy (headline + bar chart + paired scatter)
# ============================================================================
s1 = prs.slides.add_slide(blank)
set_slide_bg(s1, BG)
add_title_block(
    s1,
    "Control Effort & Energy",
    "20 paired envs, identical seed, 1000 steps, cmd_vx = 1.0 m/s  -  every metric is paired t-tested",
)

# Stat row across the top under the title
stat_y = Inches(1.40)
stat_h = Inches(0.95)
card_w = Inches(2.20)
gap_x  = Inches(0.18)
total_w = 4 * card_w + 3 * gap_x
start_x = (SLIDE_W - total_w) / 2

add_stat_card(s1, start_x + 0*(card_w+gap_x), stat_y, card_w, stat_h,
              "-4.8%", "Jerk RMS  (20/20 wins)", GREEN)
add_stat_card(s1, start_x + 1*(card_w+gap_x), stat_y, card_w, stat_h,
              "-4.1%", "Mean power  (20/20 wins)", GREEN)
add_stat_card(s1, start_x + 2*(card_w+gap_x), stat_y, card_w, stat_h,
              "-4.1%", "Total energy  (20/20 wins)", GREEN)
add_stat_card(s1, start_x + 3*(card_w+gap_x), stat_y, card_w, stat_h,
              "p < 1e-12", "All paired t-tests", ACCENT2, value_size=22)

# Bottom: bar chart (left) + paired scatter (right)
plot_y = Inches(2.50)
plot_h = Inches(2.65)
left_x = Inches(0.35)
right_x = Inches(5.25)
plot_w = Inches(4.50)

add_rect(s1, left_x - Inches(0.05),  plot_y - Inches(0.05), plot_w + Inches(0.10), plot_h + Inches(0.10), CARD)
s1.shapes.add_picture(rf"{PLOTS}\effort_bars.png", left_x, plot_y, width=plot_w)

add_rect(s1, right_x - Inches(0.05), plot_y - Inches(0.05), plot_w + Inches(0.10), plot_h + Inches(0.10), CARD)
s1.shapes.add_picture(rf"{PLOTS}\effort_paired.png", right_x, plot_y, width=plot_w)

add_footer(s1, "20 envs  |  seed=42  |  flow_model_adaptive.pt  |  t_end=0.3")


# ============================================================================
# Slide 2: Power vs time
# ============================================================================
s2 = prs.slides.add_slide(blank)
set_slide_bg(s2, BG)
add_title_block(
    s2,
    "Where the Energy Savings Come From",
    "Smoothed actions reduce instantaneous mechanical power throughout the rollout",
)

plot_y = Inches(1.50)
plot_w = Inches(6.30)
plot_x = (SLIDE_W - plot_w) / 2
plot_h = Inches(3.00)
add_rect(s2, plot_x - Inches(0.05), plot_y - Inches(0.05), plot_w + Inches(0.10), plot_h + Inches(0.10), CARD)
s2.shapes.add_picture(rf"{PLOTS}\power_vs_time.png", plot_x, plot_y, width=plot_w, height=plot_h)

# Caption strip
add_text(s2, Inches(0.5), Inches(4.70), Inches(9.0), Inches(0.55),
         "Same torque magnitude (RMS unchanged), but lower joint velocity * torque correlation  ->  "
         "less wasted mechanical work. Flow smoothing damps the high-frequency action chatter that costs energy.",
         size=12, color=TEXT, font="Calibri", align="center")

add_footer(s2, "Power = sum_j |tau_j * qdot_j|, env-mean across 20 seeds")


# ============================================================================
# Slide 3: Push recovery
# ============================================================================
s3 = prs.slides.add_slide(blank)
set_slide_bg(s3, BG)
add_title_block(
    s3,
    "Push Recovery: Smoothed Survives More Often",
    "80 paired envs, lateral impulse 200 ms wide, magnitudes swept 50-500 N",
)

# Stat row
stat_y = Inches(1.40)
stat_h = Inches(0.95)
card_w = Inches(2.20)
gap_x  = Inches(0.18)
total_w = 4 * card_w + 3 * gap_x
start_x = (SLIDE_W - total_w) / 2

add_stat_card(s3, start_x + 0*(card_w+gap_x), stat_y, card_w, stat_h,
              "34%", "Smoothed fall rate", GREEN)
add_stat_card(s3, start_x + 1*(card_w+gap_x), stat_y, card_w, stat_h,
              "50%", "Raw fall rate", RED)
add_stat_card(s3, start_x + 2*(card_w+gap_x), stat_y, card_w, stat_h,
              "p = 0.015", "McNemar (b=19, c=6)", ACCENT2, value_size=22)
add_stat_card(s3, start_x + 3*(card_w+gap_x), stat_y, card_w, stat_h,
              "+33 ms", "Avg time-to-fall (both fell)", ACCENT2)

plot_y = Inches(2.50)
plot_h = Inches(2.65)
left_x = Inches(0.35)
right_x = Inches(5.25)
plot_w = Inches(4.50)

add_rect(s3, left_x - Inches(0.05),  plot_y - Inches(0.05), plot_w + Inches(0.10), plot_h + Inches(0.10), CARD)
s3.shapes.add_picture(rf"{PLOTS}\push_fall_rate.png", left_x, plot_y, width=plot_w)

add_rect(s3, right_x - Inches(0.05), plot_y - Inches(0.05), plot_w + Inches(0.10), plot_h + Inches(0.10), CARD)
s3.shapes.add_picture(rf"{PLOTS}\push_mcnemar.png", right_x, plot_y, width=plot_w)

add_footer(s3, "Falls = early termination (excludes timeout)  |  paired by env index, identical seed=42")


# ============================================================================
# Slide 4: Push videos (optional, only if MP4s exist)
# ============================================================================
def find_mp4(folder):
    if not os.path.isdir(folder): return None
    for f in os.listdir(folder):
        if f.lower().endswith(".mp4"):
            return os.path.join(folder, f)
    return None

raw_mp4 = find_mp4(os.path.join(PUSH_VID_DIR, "raw_350"))
smo_mp4 = find_mp4(os.path.join(PUSH_VID_DIR, "smoothed_350"))

if raw_mp4 and smo_mp4:
    s4 = prs.slides.add_slide(blank)
    set_slide_bg(s4, BG)
    add_title_block(
        s4,
        "Push Recovery: Side-by-Side Example  (F = 350 N)",
        "Same seed, same 350 N impulse  -  smoothed survives +300 ms longer (raw 2.18 s -> smo 2.48 s)",
    )

    vid_y = Inches(1.55)
    vid_w = Inches(4.50)
    vid_h = Inches(2.85)
    label_h = Inches(0.40)
    gap_x = Inches(0.30)
    total_w = 2 * vid_w + gap_x
    start_x = (SLIDE_W - total_w) / 2

    # Left: raw
    lx = start_x
    add_text(s4, lx, vid_y, vid_w, label_h, "Raw PPO  -  fell at t = 2.18 s",
             size=16, color=RED, bold=True, font="Georgia", align="center")
    s4.shapes.add_movie(raw_mp4, lx, vid_y + label_h, vid_w, vid_h,
                        mime_type="video/mp4")

    # Right: smoothed
    rx = start_x + vid_w + gap_x
    add_text(s4, rx, vid_y, vid_w, label_h, "Smoothed  -  fell at t = 2.48 s",
             size=16, color=GREEN, bold=True, font="Georgia", align="center")
    s4.shapes.add_movie(smo_mp4, rx, vid_y + label_h, vid_w, vid_h,
                        mime_type="video/mp4")

    add_text(s4, Inches(0.5), Inches(5.00), Inches(9.0), Inches(0.35),
             "Click each video to play. Lateral 350 N impulse for 200 ms at t = 1.6 s. "
             "Both eventually fall at this severe magnitude, but the smoothed policy resists 300 ms longer.",
             size=10, color=MUTED, font="Calibri", align="center")
    add_footer(s4, "push_video_eval.py | seed=42 | F=350 N | push_dur=10 steps")
    n_added = 4
    print("Added 4 slides (incl. video slide)")
else:
    n_added = 3
    print("Added 3 slides (videos not found, will add slide later if needed)")

# Move conclusion to the very end
final_count = len(prs.slides)
move_slide(prs, conclusion_old_index, final_count - 1)
print(f"Final slide count: {final_count}")

prs.save(DST)
print(f"Saved to {DST}")
