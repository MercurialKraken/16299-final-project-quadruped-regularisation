"""Insert two A-to-B race slides into the flow matching deck, matching its dark theme."""
import copy
import io
import sys

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8", errors="replace")

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

# Start fresh from the backup so re-runs are deterministic
SRC = r"A:\IsaacLab\flow_matching_presentation.backup.pptx"
DST = r"A:\IsaacLab\flow_matching_presentation.pptx"
PLOTS = r"A:\IsaacLab\ab_race_plots_dark"

# Deck palette
BG      = RGBColor(0x0F, 0x17, 0x2A)  # slate-900
CARD    = RGBColor(0x1E, 0x29, 0x3B)  # slate-800
TEXT    = RGBColor(0xF8, 0xFA, 0xFC)  # slate-50
MUTED   = RGBColor(0x94, 0xA3, 0xB8)  # slate-400
ACCENT  = RGBColor(0x0E, 0xA5, 0xE9)  # sky-500
ACCENT2 = RGBColor(0x38, 0xBD, 0xF8)  # sky-400
GREEN   = RGBColor(0x22, 0xC5, 0x5E)  # green-500
RED     = RGBColor(0xEF, 0x44, 0x44)  # red-500

SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)


def set_slide_bg(slide, rgb):
    """Paint the slide background with a solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def add_rect(slide, x, y, w, h, fill_rgb, line_rgb=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_rgb
    if line_rgb is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_rgb
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, size=14, color=TEXT, bold=False, font="Calibri", align="left"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    # Ensure first paragraph
    p = tf.paragraphs[0]
    from pptx.enum.text import PP_ALIGN
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def add_stat_card(slide, x, y, w, h, value, label, value_color=GREEN):
    add_rect(slide, x, y, w, h, CARD)
    # Value (big number)
    add_text(slide, x, y + Emu(int(h * 0.12)), w, Emu(int(h * 0.55)),
             value, size=36, color=value_color, bold=True, font="Georgia", align="center")
    # Label
    add_text(slide, x, y + Emu(int(h * 0.68)), w, Emu(int(h * 0.25)),
             label, size=11, color=MUTED, font="Calibri", align="center")


def add_title_block(slide, title, subtitle):
    # Slide title
    add_text(slide, Inches(0.5), Inches(0.30), Inches(9.0), Inches(0.65),
             title, size=30, color=TEXT, bold=True, font="Georgia")
    # Subtitle (muted)
    add_text(slide, Inches(0.5), Inches(0.95), Inches(9.0), Inches(0.35),
             subtitle, size=13, color=MUTED, font="Calibri")
    # Top accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Emu(55000))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()


def add_footer(slide, text):
    add_text(slide, Inches(0.5), Inches(5.35), Inches(9.0), Inches(0.25),
             text, size=9, color=MUTED, font="Calibri", align="right")


def move_slide(prs, old_index, new_index):
    """Reorder slides inside the <p:sldIdLst>."""
    xml_slides = prs.slides._sldIdLst  # parent element
    slides = list(xml_slides)
    xml_slides.remove(slides[old_index])
    xml_slides.insert(new_index, slides[old_index])


prs = Presentation(SRC)
n_before = len(prs.slides)
print(f"Original slides: {n_before}")

print(f"Available layouts: {[l.name for l in prs.slide_layouts]}")
# Use the first (and likely only) layout — matches existing slides' 'DEFAULT'
blank = prs.slide_layouts[0]

# ============== Slide A: A-to-B Race — quantitative ==============
sA = prs.slides.add_slide(blank)
set_slide_bg(sA, BG)
add_title_block(
    sA,
    "A-to-B Race: Flat-Terrain Head-to-Head",
    "500 steps, cmd_vx = 1.0 m/s — flow-smoothed PPO beats raw in every metric",
)

# Left half: 2x2 stat cards
card_w = Inches(2.15)
card_h = Inches(1.30)
left_x = Inches(0.5)
gap_x  = Inches(0.25)
row1_y = Inches(1.55)
row2_y = Inches(3.05)

add_stat_card(sA, left_x,                                row1_y, card_w, card_h,
              "+14 cm", "Lead distance over raw", GREEN)
add_stat_card(sA, left_x + card_w + gap_x,               row1_y, card_w, card_h,
              "+2.7%", "Mean forward velocity", ACCENT2)
add_stat_card(sA, left_x,                                row2_y, card_w, card_h,
              "+6.1%", "Peak forward velocity", ACCENT2)
add_stat_card(sA, left_x + card_w + gap_x,               row2_y, card_w, card_h,
              "-2.9%", "Path length (straighter)", GREEN)

# Right half: progress_vs_time plot inside a card
plot_x = Inches(5.35)
plot_y = Inches(1.45)
plot_w = Inches(4.25)
card_p = add_rect(sA, plot_x - Inches(0.05), plot_y - Inches(0.05),
                   plot_w + Inches(0.10), Inches(3.10), CARD)
sA.shapes.add_picture(rf"{PLOTS}\progress_vs_time.png", plot_x, plot_y, width=plot_w)

# Caption under plot
add_text(sA, plot_x - Inches(0.05), plot_y + Inches(2.85), plot_w + Inches(0.10), Inches(0.3),
         "Forward progress along x — smoothed trajectory pulls ahead around t = 2s and stays ahead",
         size=10, color=MUTED, font="Calibri", align="center")

add_footer(sA, "Raw vs flow-smoothed (t_end=0.3, adaptive Butterworth)  |  checkpoint: model_299, flat PPO")

# ============== Slide B: A-to-B Race — trajectory + actions ==============
sB = prs.slides.add_slide(blank)
set_slide_bg(sB, BG)
add_title_block(
    sB,
    "Race Detail: Trajectory, Velocity, and Action Smoothness",
    "Smoothed policy walks a straighter line, tracks cmd more tightly, and halves per-step jitter",
)

# 3-column layout: all plots now share ~1.15 aspect ratio so they line up.
# Each picture: width=2.90" → height ≈ 2.90/1.154 = 2.51"
p_w = Inches(2.90)
p_gap = Inches(0.15)
total_w = 3 * p_w + 2 * p_gap  # 8.70"
start_x = (SLIDE_W - total_w) / 2
row_y = Inches(1.50)
row_h = Inches(3.20)  # card height: fits 2.51" picture + 0.30" caption + padding

plots = [
    ("path_xy.png",     "Top-down XY path — smoothed line is straighter"),
    ("vx_vs_time.png",  "Forward velocity tracking (cmd = 1.0 m/s)"),
    ("action_rate.png", "Per-step action change (jitter proxy)"),
]

for i, (fn, cap) in enumerate(plots):
    x = start_x + i * (p_w + p_gap)
    add_rect(sB, x - Inches(0.05), row_y - Inches(0.05),
             p_w + Inches(0.10), row_h + Inches(0.10), CARD)
    sB.shapes.add_picture(rf"{PLOTS}\{fn}", x, row_y, width=p_w)
    # Caption at bottom of card
    add_text(sB, x, row_y + row_h - Inches(0.30), p_w, Inches(0.30),
             cap, size=10, color=MUTED, font="Calibri", align="center")

# Bottom summary strip
strip_y = Inches(4.95)
add_text(sB, Inches(0.5), strip_y, Inches(9.0), Inches(0.25),
         "Smoothed:  final_x = 5.76 m  |  mean v_x = 0.578 m/s  |  peak v_x = 0.87 m/s  |  path = 10.12 m   "
         "\u2014   Raw:  5.62 m  |  0.563 m/s  |  0.82 m/s  |  10.43 m",
         size=10, color=TEXT, font="Calibri", align="center")

add_footer(sB, "Videos:  smoothed_t0.0  (raw)   &   smoothed_t0.3  (flow-smoothed)   in  logs\\rsl_rl\\unitree_go1_flat\\2026-04-06_12-42-26\\videos")

# Reorder: place new slides immediately BEFORE the Conclusion slide (originally last)
# Indices at this point: 0..n_before-1 are original slides; the two new slides are n_before and n_before+1.
# The original conclusion slide is at index n_before-1.
#
# Target order:
#   original slides 0..n_before-2 (everything except the original conclusion)
#   new slide A (currently index n_before)
#   new slide B (currently index n_before+1)
#   original conclusion (currently index n_before-1)
#
# Move new slide A to position n_before-1 (just before original conclusion)
# Move new slide B to position n_before   (just before original conclusion, after A)
# After moving A: A lives at index n_before-1, conclusion is pushed to n_before.
# After moving B (originally at new_index n_before+1 but indices shifted) — easiest: move conclusion to very end.

# Simpler: just move the original conclusion slide to the very end.
conclusion_old_index = n_before - 1
# After appending 2 slides, conclusion is still at index n_before-1, and the last index is n_before+1.
move_slide(prs, conclusion_old_index, n_before + 1)

print(f"Final slide count: {len(prs.slides)}")
prs.save(DST)
print(f"Saved to {DST}")
